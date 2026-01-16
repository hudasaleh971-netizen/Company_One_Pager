"""
PPT Slide Generator
====================
Generates PowerPoint slides from SlideData JSON using Template.pptx.
Uses python-pptx to replace placeholders and create tables.
"""

import os
import re
from pathlib import Path
from typing import Dict, Any, Optional
from loguru import logger

from app.models.ppt_models import SlideData, ManagementMember, Shareholder


# Project paths
PROJECT_ROOT = Path(__file__).parent.parent
TEMPLATE_PATH = PROJECT_ROOT / "Template.pptx"
OUTPUT_DIR = PROJECT_ROOT / "outputs"


def _get_pptx_module():
    """Lazy import of python-pptx to avoid import-time errors."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        return Presentation, True
    except ImportError:
        logger.warning("python-pptx not installed. Run: pip install python-pptx")
        return None, False


def replace_placeholder_in_shape(shape, placeholder: str, value: str) -> bool:
    """Replace a placeholder in a shape's text. Returns True if replaced."""
    if not hasattr(shape, "text_frame"):
        return False
    
    replaced = False
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if placeholder in run.text:
                run.text = run.text.replace(placeholder, value)
                replaced = True
    return replaced


def replace_placeholders_in_slide(slide, replacements: Dict[str, str]):
    """Replace all placeholders in a slide."""
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for placeholder, value in replacements.items():
                replace_placeholder_in_shape(shape, placeholder, value)
        
        # Handle tables
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for placeholder, value in replacements.items():
                        if placeholder in cell.text:
                            cell.text = cell.text.replace(placeholder, value)


def create_management_table_text(members: list) -> str:
    """Format management team data as text for placeholder replacement."""
    if not members:
        return "No management data available"
    
    lines = []
    for member in members:
        if isinstance(member, dict):
            name = member.get("name", "")
            position = member.get("position", "")
            bio = member.get("bio", "")
        else:
            name = member.name
            position = member.position
            bio = member.bio
        
        lines.append(f"{name}\n{position}\n{bio}")
        lines.append("-" * 50)
    
    return "\n".join(lines[:-1])  # Remove last separator


def create_shareholders_table_text(shareholders: list) -> str:
    """Format shareholders data as text for placeholder replacement."""
    if not shareholders:
        return "No shareholder data available"
    
    lines = []
    for sh in shareholders:
        if isinstance(sh, dict):
            name = sh.get("name", "")
            pct = sh.get("ownership_percentage", "n/d")
        else:
            name = sh.name
            pct = sh.ownership_percentage
        
        lines.append(f"{name}: {pct}")
    
    return "\n".join(lines)


def generate_slides_from_data(slide_data: Dict[str, Any], output_filename: str = None) -> Optional[str]:
    """
    Generate PowerPoint slides from SlideData dictionary.
    
    Args:
        slide_data: Dictionary matching SlideData model fields
        output_filename: Optional custom output filename
        
    Returns:
        Path to generated file, or None if generation failed
    """
    # Lazy import python-pptx
    Presentation, pptx_available = _get_pptx_module()
    
    if not pptx_available:
        logger.error("python-pptx not installed. Cannot generate slides.")
        return None
    
    if not TEMPLATE_PATH.exists():
        logger.error(f"Template not found: {TEMPLATE_PATH}")
        return None
    
    # Ensure output directory exists
    OUTPUT_DIR.mkdir(exist_ok=True)
    
    # Generate output filename
    if output_filename is None:
        company_name = slide_data.get("COMPANY_NAME", "Company").replace(" ", "_")
        output_filename = f"{company_name}_OnePager.pptx"
    
    output_path = OUTPUT_DIR / output_filename
    
    logger.info(f"📊 Generating PPT from template: {TEMPLATE_PATH}")
    logger.info(f"📁 Output: {output_path}")
    
    try:
        # Load template
        prs = Presentation(str(TEMPLATE_PATH))
        
        # Build replacements dictionary
        replacements = {
            "{{COMPANY_NAME}}": slide_data.get("COMPANY_NAME", ""),
            "{{COUNTRY}}": slide_data.get("COUNTRY", ""),
            "{{BACKGROUND_SUMMARY}}": slide_data.get("BACKGROUND_SUMMARY", ""),
            "{{KEY_PRODUCTS}}": slide_data.get("KEY_PRODUCTS", ""),
            "{{Unit}}": slide_data.get("Unit", "USDm"),
            "{{year}}": slide_data.get("year", "2025"),
            "{{BORROWERS_VALUE}}": slide_data.get("BORROWERS_VALUE", "n/d"),
            "{{EMPLOYEES_VALUE}}": slide_data.get("EMPLOYEES_VALUE", "n/d"),
            "{{OUTSTANDING_VALUE}}": slide_data.get("OUTSTANDING_VALUE", "n/d"),
            "{{PAR_VALUE}}": slide_data.get("PAR_VALUE", "n/d"),
            "{{DISBURSALS_VALUE}}": slide_data.get("DISBURSALS_VALUE", "n/d"),
            "{{EQUITY_VALUE}}": slide_data.get("EQUITY_VALUE", "n/d"),
            "{{NET_INCOME_VALUE}}": slide_data.get("NET_INCOME_VALUE", "n/d"),
            "{{CREDIT_RATING_VALUE}}": slide_data.get("CREDIT_RATING_VALUE", "n/d"),
            "{{TABLE_MANAGEMENT}}": create_management_table_text(
                slide_data.get("TABLE_MANAGEMENT", [])
            ),
            "{{TABLE_SHAREHOLDERS}}": create_shareholders_table_text(
                slide_data.get("TABLE_SHAREHOLDERS", [])
            ),
        }
        
        # Replace placeholders in all slides
        for slide in prs.slides:
            replace_placeholders_in_slide(slide, replacements)
        
        # Save output
        prs.save(str(output_path))
        logger.info(f"✅ PPT generated successfully: {output_path}")
        
        return str(output_path)
        
    except Exception as e:
        logger.error(f"❌ Error generating PPT: {e}")
        import traceback
        logger.error(traceback.format_exc())
        return None


def generate_ppt_callback(callback_context):
    """
    After-agent callback for LoopAgent/SequentialAgent.
    Generates PPT slides when the agent completes.
    
    This callback reads 'slide_data' from session state and generates
    a PowerPoint file using Template.pptx.
    
    Returns:
        None (callback does not modify agent response)
    """
    logger.info("📋 [Callback] PPT Generation triggered")
    
    state = callback_context.state
    slide_data = state.get("slide_data")
    
    if not slide_data:
        logger.warning("⚠️ No slide_data in session state. Skipping PPT generation.")
        return None
    
    # Parse if it's a string (JSON)
    if isinstance(slide_data, str):
        import json
        try:
            slide_data = json.loads(slide_data)
        except json.JSONDecodeError as e:
            logger.error(f"❌ Failed to parse slide_data JSON: {e}")
            return None
    
    # Generate PPT
    output_path = generate_slides_from_data(slide_data)
    
    if output_path:
        # Store output path in state for API to retrieve
        state["ppt_output_path"] = output_path
        logger.info(f"✅ PPT saved to: {output_path}")
    
    return None  # Don't modify agent response
